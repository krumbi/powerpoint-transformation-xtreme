import argparse
import collections.abc
import pptx
import pathlib
import logging as log
import sys
import openai
import json
import re

language_map = {"en": "English", "de": "German"}

parser = argparse.ArgumentParser()
parser.add_argument('input', metavar='INPUT', type=str, help='the input file')
parser.add_argument('key', metavar='KEY', type=str, help='the OpenAI API key')
parser.add_argument('-o', '--output', metavar='OUTPUT', default=pathlib.Path("output"), type=pathlib.Path, help='the output directory')
parser.add_argument('-m', '--model', metavar='MODEL', default="gpt-3.5-turbo-16k", type=str, help='the OpenAI model')
parser.add_argument('-t', '--temperature', metavar='TEMPERATURE', default=0.3, type=float, help='the OpenAI temperature')
parser.add_argument('-l', '--language', metavar='LANGUAGE', default="en", choices=['en', 'de'], help='the language of the output text')
parser.add_argument('-p', '--prompt', metavar='PROMPT', default="", type=str, help='the prompt extension (e.g. "As a lecture")')
parser.add_argument('-d', '--debug', action='store_true', help='enable debug logging')
args = parser.parse_args()

log.basicConfig(level=log.INFO if not args.debug else log.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

log.info("Starting text-to-text")

for var, val in vars(args).items():
    log.info(f"Input: {var} = {val}")

if args.debug:
    log.debug("Debug mode enabled")

if not (input_file := pathlib.Path(args.input)).is_file():
    log.error(f"Input {input_file} does not exist")
    sys.exit(1)
if input_file.suffix != ".pptx":
    log.error(f"Input {input_file} is not a .pptx file")
    sys.exit(1)

if not args.output.exists():
    log.error(f"Output directory {args.output} does not exist")
    sys.exit(1)

prs = pptx.Presentation(args.input)

raw_text = []

for i, s in enumerate(prs.slides, start=1):
    text = []
    for shape in s.shapes:
        if not shape.has_text_frame or len(shape.text) == 0:
            continue
        if "Title" in shape.name:
            text.append(f"Title: {shape.text}")
            continue
        
        for p in shape.text_frame.paragraphs:
            if len(p.runs) > 0:
                current_text = ""
                for r in p.runs:
                    color_type = str(r.font.color.type)
                    if "SCHEME" in color_type and "BACKGROUND" in str(r.font.color.theme_color):
                        pass
                    elif color_type == "None":
                        current_text += r.text
                    else:
                        current_text += r.text
                if len(current_text) == 0:
                    continue
                text.append(f"{'  ' * p.level}- {current_text}")
            else:    
                if len(p.text) == 0:
                    continue
                text.append(f"{'  ' * p.level}- {p.text}")
        
    raw_text.append("\n".join(text))

    log.info(f"Processed Slide {i}")

log.info(f"Processing of slides ({len(prs.slides)}) finished")

debug_folder = args.output.joinpath("debug")

if args.debug:
    debug_folder.mkdir(parents=True, exist_ok=True)

    raw_text_folder = debug_folder.joinpath("raw_text")
    raw_text_folder.mkdir(parents=True, exist_ok=True)

    log.debug(f"Writing raw text to {raw_text_folder.absolute()}...")

    for i, text in enumerate(raw_text, start=1):
        with open(raw_text_folder.joinpath(f"slide_{i}.txt"), "w", encoding="utf-8") as f:
                f.write(text)

    log.debug("Writing raw text finished")

log.info("Generating ChatGPT prompt...")

chatgpt_prompt = \
    f"""Can you generate me a text for a spoken presentation based on the following information extracted from presentation slides?
Please keep the separation in slides. It should be a bit more formal and structured. Please add some information where it is useful.
Please generate the text in {language_map[args.language]}. {args.prompt}
Thank you very much!"""

chatgpt_prompt += "\n\n"

for i, text in enumerate(raw_text, start=1):
    chatgpt_prompt += f"Slide {i}:\n{text}\n\n"

chatgpt_prompt.strip()

log.log(log.INFO, "Generating ChatGPT prompt finished")


if args.debug:
    log.debug(f"Writing ChatGPT prompt to {args.output.absolute()}")

    with open(debug_folder.joinpath("prompt.txt"), "w", encoding="utf-8") as f:
        f.write(chatgpt_prompt)
    
    log.debug("Writing ChatGPT prompt finished")

if not args.debug:
    openai.api_key = args.key

    log.info("Requesting ChatGPT completion...")
    #completion = openai.Completion.create(model=args.model, prompt=chatgpt_prompt, temperature=args.temperature, max_tokens=10000)
    completion = openai.ChatCompletion.create(model=args.model, messages=[{"role": "user", "content": chatgpt_prompt}], temperature=args.temperature, max_tokens=10000)

    with open(args.output.joinpath("chat.json"), "w", encoding="utf-8") as f:
        json.dump(completion, f, indent=4)

    log.info("ChatGPT completion finished")
else:
    from types import SimpleNamespace
    log.warning("Skipping ChatGPT request. Reading from file.")
    with open(args.output.joinpath("chat.json"), "r", encoding="utf-8") as f:
        completion = json.load(f, object_hook=lambda d: SimpleNamespace(**d))

if len(completion.choices) > 1:
    log.warning("More than one choice returned by ChatGPT")

text : str = completion.choices[0].message.content

log.info("Splitting text...")

split_pattern = re.compile(r"Slide (\d+):")
splitted_text = [part for part in map(str.strip, split_pattern.split(text)) if len(part) > 0]

log.info("Splitted text into slides")

log.info("Writing slides to files...")

for i in range(1, len(splitted_text), 2):
    with open(args.output.joinpath(f"slide_{splitted_text[i - 1]}.txt"), "w", encoding="utf-8") as f:
        f.write(splitted_text[i])

log.info("Writing slides to files finished")

log.info("Text-to-text finished")
